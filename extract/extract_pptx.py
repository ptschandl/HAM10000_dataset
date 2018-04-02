#!/usr/bin/env python3

"""
Extract images with corresponding label from PPTX-presentations.
(c) Philipp Tschandl, 2018
You may re-use this code under the CC BY-NC 4.0 license.

Instructions:
- Run "python extract_pptx.py"
- Presentations:
    - ...should span a single calender-year and have it as
    a string within the file-name (e.g. "sample_2001.pptx")
    - ...should be in the "presentations"-subfolder
    - ...have to be in .pptx format. If in .ppt format,
one can transform them in batch by the libreoffice CLI:

libreoffice --headless --invisible --convert-to pptx *.ppt



Output:
- Debug information during extraction is stored in a .log file
- Images will be stored in images/ in the following format:
    - ID_YEAR_SLIDENr.jpg
        - ID = Text-Field found on same slide
        - YEAR = "20**" identifyer of .pptx-filename
        - SLIDENr = Slide ID (Note: this does not start at 1)

"""

import pptx # pip install python-pptx
from pptx import Presentation
from PIL import Image
from tqdm import tqdm
import io
import re
import os
import glob

import datetime
date = datetime.datetime.now().strftime("%Y-%m-%d")

import logging
logging.basicConfig(filename='extract_pptx_' + date + '.log', level=logging.DEBUG)

################################################################################################
# GET pptx files
################################################################################################

filenames = glob.glob("presentations/*.pptx")
years = [re.search(r"(20|19)\d\d", s).group() for s in filenames]
print("Number of presentation files: {}".format(len(filenames)))
print("Years covered: {}".format(sorted(set(years))))

all_lesions = {}

################################################################################################
# EXTRACT Information and Images
################################################################################################

for f in tqdm(filenames, desc="PPTX-Files"):

    year = re.search(r"(20|19)\d\d", f).group()

    try:
        prs = Presentation(f)
    except:
        logging.debug("Error reading presentation:" + f)
        continue

    nr_slides = len(prs.slides)
    data = {}
    im = ""
    label = ""

    for slide in tqdm(prs.slides, desc="Read Slides", leave=False):

        for shape in slide.shapes:

            # Find the image shape
            if type(shape) != pptx.shapes.autoshape.Shape:  # Cannot verify image-shape on shapes.autoshapes
                if shape.shape_type == 13:  # Check if shape_type is an image
                    try:
                        im = shape.image.blob
                    except:
                        im = ""
                        logging.debug('Image extraction error in file: ' + f + " on slide: " + str(slide.slide_id))
                    continue

            if not shape.has_text_frame:
                continue

            # Find the text shape
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    label = run.text

        # Use data only if image AND label existed
        if im is not None and im != "" and label is not None and label != "":
            data[slide.slide_id] = [label, im]

        # Save datapoint to dict for debugging purposes
        all_lesions[str(label).strip() + "_" + str(year) + "_" + str(slide.slide_id)] = [f, slide.slide_id,
                                                                                         slide.slide_id - 256]

    # Store all found images
    os.makedirs("./images", exist_ok=True)
    for index, (lab, im) in tqdm(data.items(), desc="Save Images", leave=False):

        # Omit slides with no label or only textlabels == redundant security to above feeding control mechanism
        if re.search(r'\d+', str(lab)) is not None and im is not None and im is not "":
            filename = re.search(r'\d+', str(lab).strip()).group() + "_" + str(year) + "_" + str(index) + ".jpg"
            img = Image.open(io.BytesIO(im))

            if img.format == "PNG":
                img = img.convert("RGB")
                img.save(os.path.join("images", filename))
                # Verify .save() call was successful
                if not os.path.isfile(os.path.join("images", filename)):
                    img.save(os.path.join("images", filename))

            else:
                img.save(os.path.join("images", filename))

################################################################################################
# CLEANUP images without proper information
################################################################################################

trueimages = [key.replace(" ", "_") + ".jpg" for key, values in all_lesions.items()]
presentimages = os.listdir("./images")
removed = 0

for testimage in presentimages:
    if testimage not in trueimages:
        os.remove("./images/" + testimage)
        removed += 1

print("Removed {} files in cleanup.".format(removed))
print("Extracted images: {}.".format(len(os.listdir("./images"))))
